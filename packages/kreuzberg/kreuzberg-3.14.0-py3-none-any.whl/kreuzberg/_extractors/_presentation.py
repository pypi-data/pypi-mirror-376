from __future__ import annotations

import re
from contextlib import suppress
from html import escape
from io import BytesIO
from pathlib import Path
from typing import TYPE_CHECKING, ClassVar

import pptx
from anyio import Path as AsyncPath
from pptx.enum.shapes import MSO_SHAPE_TYPE

from kreuzberg._extractors._base import Extractor
from kreuzberg._mime_types import MARKDOWN_MIME_TYPE, POWER_POINT_MIME_TYPE
from kreuzberg._types import ExtractionResult
from kreuzberg._utils._string import normalize_spaces

if TYPE_CHECKING:  # pragma: no cover
    from pptx.presentation import Presentation

    from kreuzberg._types import Metadata

_NON_WORD_PATTERN = re.compile(r"\W")


class PresentationExtractor(Extractor):
    SUPPORTED_MIME_TYPES: ClassVar[set[str]] = {POWER_POINT_MIME_TYPE}

    async def extract_bytes_async(self, content: bytes) -> ExtractionResult:
        return self._extract_pptx(content)

    async def extract_path_async(self, path: Path) -> ExtractionResult:
        content = await AsyncPath(path).read_bytes()
        return self._extract_pptx(content)

    def extract_bytes_sync(self, content: bytes) -> ExtractionResult:
        return self._extract_pptx(content)

    def extract_path_sync(self, path: Path) -> ExtractionResult:
        content = Path(path).read_bytes()
        return self._extract_pptx(content)

    def _extract_pptx(self, file_contents: bytes) -> ExtractionResult:
        md_content = ""
        presentation = pptx.Presentation(BytesIO(file_contents))

        for index, slide in enumerate(presentation.slides):
            md_content += f"\n\n<!-- Slide number: {index + 1} -->\n"

            title = None
            if hasattr(slide.shapes, "title"):
                title = slide.shapes.title

            for shape in slide.shapes:
                if not hasattr(shape, "shape_type"):
                    continue

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (
                    shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image")
                ):
                    alt_text = ""
                    with suppress(AttributeError):
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")  # noqa: SLF001

                    filename = _NON_WORD_PATTERN.sub("", shape.name) + ".jpg"
                    md_content += f"\n![{alt_text if alt_text else shape.name}]({filename})\n"

                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    html_table = "<table>"
                    first_row = True

                    for row in shape.table.rows:
                        html_table += "<tr>"

                        for cell in row.cells:
                            tag = "th" if first_row else "td"
                            html_table += f"<{tag}>{escape(cell.text)}</{tag}>"

                        html_table += "</tr>"
                        first_row = False

                    html_table += "</table>"
                    md_content += "\n" + html_table + "\n"

                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

            md_content = md_content.strip()
            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame

                if notes_frame is not None:  # pragma: no branch
                    md_content += notes_frame.text

                md_content = md_content.strip()

        result = ExtractionResult(
            content=normalize_spaces(md_content),
            mime_type=MARKDOWN_MIME_TYPE,
            metadata=self._extract_presentation_metadata(presentation),
            chunks=[],
        )

        return self._apply_quality_processing(result)

    @staticmethod
    def _extract_presentation_metadata(presentation: Presentation) -> Metadata:
        metadata: Metadata = {}

        PresentationExtractor._extract_core_properties(presentation, metadata)

        fonts = PresentationExtractor._extract_fonts(presentation)
        if fonts:
            metadata["fonts"] = list(fonts)

        PresentationExtractor._add_presentation_structure_info(presentation, metadata, fonts)

        return metadata

    @staticmethod
    def _extract_core_properties(presentation: Presentation, metadata: Metadata) -> None:
        property_mapping = [
            ("authors", "author"),
            ("comments", "comments"),
            ("status", "content_status"),
            ("created_by", "created"),
            ("identifier", "identifier"),
            ("keywords", "keywords"),
            ("modified_by", "last_modified_by"),
            ("modified_at", "modified"),
            ("version", "revision"),
            ("subject", "subject"),
            ("title", "title"),
        ]

        for metadata_key, core_property_key in property_mapping:
            if core_property := getattr(presentation.core_properties, core_property_key, None):
                metadata[metadata_key] = core_property  # type: ignore[literal-required]

        if presentation.core_properties.language:
            metadata["languages"] = [presentation.core_properties.language]

        if presentation.core_properties.category:
            metadata["categories"] = [presentation.core_properties.category]

    @staticmethod
    def _extract_fonts(presentation: Presentation) -> set[str]:
        fonts = set()
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run, "font") and run.font.name:
                            fonts.add(run.font.name)
        return fonts

    @staticmethod
    def _add_presentation_structure_info(presentation: Presentation, metadata: Metadata, fonts: set[str]) -> None:
        slide_count = len(presentation.slides)
        if slide_count == 0:
            return

        structure_info = f"Presentation with {slide_count} slide{'s' if slide_count != 1 else ''}"

        slides_with_notes = sum(1 for slide in presentation.slides if slide.has_notes_slide)
        if slides_with_notes > 0:
            structure_info += f", {slides_with_notes} with notes"

        metadata["description"] = structure_info

        if "summary" not in metadata:
            summary_parts = [f"PowerPoint presentation with {slide_count} slides"]
            if slides_with_notes > 0:
                summary_parts.append(f"{slides_with_notes} slides have notes")
            if fonts:
                summary_parts.append(f"uses {len(fonts)} font{'s' if len(fonts) != 1 else ''}")

            metadata["summary"] = f"{'. '.join(summary_parts)}."
